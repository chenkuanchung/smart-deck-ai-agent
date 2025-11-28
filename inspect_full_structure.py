# inspect_full_structure.py
from pptx import Presentation
import os

def analyze_structure(filename="template.pptx"):
    if not os.path.exists(filename):
        print(f"❌ 找不到 {filename}")
        return

    prs = Presentation(filename)
    print(f"🔍 正在深度分析 {filename} 的結構...\n")
    print("="*60)

    for i, layout in enumerate(prs.slide_layouts):
        print(f"📂 版型 ID: [{i}] - 名稱: {layout.name}")
        
        if not layout.placeholders:
            print("   (此版型沒有任何可填寫的框框)")
        else:
            for shape in layout.placeholders:
                # 判斷框框類型
                p_type = shape.placeholder_format.type
                # 印出 ID 和類型
                print(f"   └── 框框 Index: [{shape.placeholder_format.idx}] - 類型: {p_type} - 預設文字: '{shape.name}'")
        
        print("-" * 60)

if __name__ == "__main__":
    analyze_structure()