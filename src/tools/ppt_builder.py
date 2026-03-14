# src/tools/ppt_builder.py
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import os

def get_layout_mapping(prs):
    """動態掃描 PPT 母片，自動找出正確的 Layout ID"""
    # 預設 Fallback (相容原本寫死的邏輯，以防母片名稱全自訂)
    mapping = {
        "title": {"id": 0, "title_idx": 0, "body_idx": 1},
        "content": {"id": 1, "title_idx": 0, "body_idx": 1},
        "section": {"id": 2, "title_idx": 0, "body_idx": 1},
        "two_column": {"id": 3, "title_idx": 0, "left_idx": 1, "right_idx": 2}
    }

    # 遍歷所有母片，根據名稱特徵自動覆蓋 ID
    for i, layout in enumerate(prs.slide_layouts):
        name = layout.name.lower()
        if "title slide" in name or "標題投影片" in name:
            mapping["title"]["id"] = i
        elif "two content" in name or "兩項內容" in name or "雙欄" in name:
            mapping["two_column"]["id"] = i
        elif "section header" in name or "章節標題" in name:
            mapping["section"]["id"] = i
        elif "title and content" in name or "標題及內容" in name:
            mapping["content"]["id"] = i

    return mapping

def set_font(paragraph, font_name="Microsoft JhengHei", size=None):
    """通用字體設定，確保中文顯示正常"""
    for run in paragraph.runs:
        run.font.name = font_name
        if size: run.font.size = size
        try:
            run.font.element.rPr.ea.attrib['typeface'] = font_name
        except: pass

def fill_text_frame(text_frame, content_items):
    """將 ContentItem 列表填入 TextFrame"""
    if not content_items: return
    text_frame.clear() 

    for i, item in enumerate(content_items):
        if hasattr(item, 'text'):
            text_val = item.text
            level_val = item.level
        elif isinstance(item, dict):
            text_val = item.get('text', '')
            level_val = item.get('level', 0)
        else:
            text_val = str(item)
            level_val = 0

        if not text_val.strip(): continue

        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
            
        p.text = text_val.strip()
        
        # 確保 level_val 一定是 int，並且限制在 0-4 之間
        try:
            safe_level = int(level_val)
        except (ValueError, TypeError):
            safe_level = 0
            
        safe_level = min(max(0, safe_level), 4) # 限制最大縮排層級為 4
        p.level = safe_level
        
        # 重新計算字體：level=0(26), level=1(22), level>=2(18)
        calc_size = 26 - (min(safe_level, 2) * 4)
        set_font(p, size=Pt(calc_size))

def create_presentation(title, slides_content, template_path="template.pptx", filename="output.pptx"):
    if os.path.exists(template_path):
        prs = Presentation(template_path)
    else:
        prs = Presentation()
        print("⚠️ 警告: 找不到 template.pptx，使用預設白底範本。")

    # 動態取得當前範本的版型配置表
    layout_config = get_layout_mapping(prs)

    for page in slides_content:
        layout_key = page.get('layout', 'content')
        config = layout_config.get(layout_key, layout_config['content'])
        
        try:
            slide_layout = prs.slide_layouts[config['id']]
        except:
            print(f"⚠️ 警告: 版型 ID {config.get('id')} 不存在，回退至預設。")
            slide_layout = prs.slide_layouts[1] 

        slide = prs.slides.add_slide(slide_layout)
        
        # ✨ 使用安全的方式獲取 Placeholder，避免 KeyError 或 IndexError
        def safe_get_placeholder(slide, idx):
            try:
                # 確保 idx 存在於 slide.placeholders 中
                if idx in [p.placeholder_format.idx for p in slide.placeholders]:
                    for p in slide.placeholders:
                        if p.placeholder_format.idx == idx:
                            return p
                # 如果找不到對應的 idx，退而求其次用陣列索引
                return slide.placeholders[idx]
            except Exception as e:
                print(f"⚠️ 警告: 找不到 Placeholder (idx={idx})，忽略此區塊。({e})")
                return None

        # A. 標題
        title_ph = safe_get_placeholder(slide, config['title_idx'])
        if title_ph and title_ph.has_text_frame:
            title_ph.text = page.get('title', '')
            set_font(title_ph.text_frame.paragraphs[0], size=Pt(32))

        # B. 內容 (雙欄與單欄分流)
        raw_items = page.get('content', [])
        if not isinstance(raw_items, list): raw_items = []

        if layout_key == 'two_column':
            left_items, right_items = [], []
            for item in raw_items:
                c = getattr(item, 'column', 0) if not isinstance(item, dict) else item.get('column', 0)
                if c == 1: right_items.append(item)
                else: left_items.append(item)

            left_ph = safe_get_placeholder(slide, config.get('left_idx', 1))
            if left_ph and left_ph.has_text_frame:
                fill_text_frame(left_ph.text_frame, left_items)
                
            right_ph = safe_get_placeholder(slide, config.get('right_idx', 2))
            if right_ph and right_ph.has_text_frame:
                fill_text_frame(right_ph.text_frame, right_items)
        else:
            body_ph = safe_get_placeholder(slide, config.get('body_idx', 1))
            if body_ph and body_ph.has_text_frame:
                fill_text_frame(body_ph.text_frame, raw_items)

        # C. 備忘稿
        if page.get('notes') and slide.has_notes_slide:
            slide.notes_slide.notes_text_frame.text = str(page.get('notes'))

    output_path = os.path.join(os.getcwd(), filename)
    prs.save(output_path)
    return output_path